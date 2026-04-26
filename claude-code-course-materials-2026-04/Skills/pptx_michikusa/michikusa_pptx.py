"""Michikusa PPTX skill helper.

Usage:
    from michikusa_pptx import Deck
    d = Deck()
    d.cover("タイトル", "サブタイトル")
    d.section("セクション区切り")
    d.bullets("見出し", ["項目1", "項目2"])
    d.twocol("見出し", ["左1","左2"], ["右1","右2"], left_head="左見出し", right_head="右見出し")
    d.save("out.pptx")

The deck is built on top of michikusa.pptx (Michikusaブランドのテンプレ):
- グラデ枠（水色→青）+ 右上ロゴ（マスター継承）
- Cover / Section / Bullets / TwoCol の4レイアウト
- 黒の太字ブレット、本文はタイトル左端に揃う
- BudouXで日本語の自然な改行
- タイトルは上部(top=0.3")に寄り、下線はタイトル直下(top=1.17")
- TwoColには専用のサブタイトルプレースホルダー (idx=10 左, idx=11 右) がある
  PowerPointで手動挿入しても5枠(タイトル/左サブ/左本文/右サブ/右本文)が見える
"""
import os
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

import budoux

_PARSER = budoux.load_default_japanese_parser()
_TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "michikusa.pptx")
_TEMPLATE_INTRO = os.path.join(os.path.dirname(os.path.abspath(__file__)), "michikusa_intro.pptx")


def _char_width(c):
    """全角換算の文字幅。日本語フォントで全角表示される記号も1.0扱い。"""
    o = ord(c)
    # CJK統合漢字・かな・全角英数・CJK記号（既存範囲）
    if (0x3000 <= o <= 0x9FFF) or (0xFF00 <= o <= 0xFFEF):
        return 1.0
    # 一般記号（—, –, …, ※ など）と矢印（→, ←, ⇒ など）も日本語組版では全角
    if (0x2010 <= o <= 0x206F) or (0x2190 <= o <= 0x21FF) or (0x25A0 <= o <= 0x25FF) or (0x2600 <= o <= 0x27BF):
        return 1.0
    return 0.55


# 寡婦回避のしきい値（全角換算）: 最後の行がこの幅以下なら前行末と結合する
ORPHAN_MIN_W = 3.0


def jp_wrap(text: str, max_chars: int = 24, orphan_min_w: float = ORPHAN_MIN_W) -> str:
    """BudouXで分かち書きしたうえで、max_chars(全角換算)を超えそうなら折り返す。

    寡婦回避: 折り返した結果、末尾行の幅が orphan_min_w 以下なら直前の行に吸収する。
    これによって「プロンプト…大事なの / は / 言語化力。」のような孤立1〜2文字行を防ぐ。
    """
    out_lines = []
    for raw_line in text.split("\n"):
        if not raw_line:
            out_lines.append("")
            continue
        phrases = _PARSER.parse(raw_line)
        cur, cur_w = "", 0.0
        line_buf = []
        for ph in phrases:
            ph_w = sum(_char_width(c) for c in ph)
            if cur and (cur_w + ph_w) > max_chars:
                line_buf.append(cur)
                cur, cur_w = ph, ph_w
            else:
                cur += ph
                cur_w += ph_w
        if cur:
            line_buf.append(cur)

        # 寡婦回避: 末尾行が極端に短ければ前行末に吸収
        while len(line_buf) >= 2:
            last_w = sum(_char_width(c) for c in line_buf[-1])
            if last_w <= orphan_min_w:
                line_buf[-2] = line_buf[-2] + line_buf[-1]
                line_buf.pop()
            else:
                break

        out_lines.extend(line_buf)
    return "\n".join(out_lines)


# Default max_chars per layout font size (全角換算).
# 実プレースホルダ幅とフォントサイズに対して安全側に設定。大きい値にすると
# PowerPoint側で自動折り返しが発生し、孤立行の原因になる。
WRAP = {
    "cover_title": 12,
    "cover_subtitle": 24,
    "section": 16,
    "bullets_title": 22,
    "bullets_body": 28,
    "twocol_title": 22,
    "twocol_head": 16,
    "twocol_body": 16,
}

# Font size (pt) used for the TwoCol column sub-title (left_head/right_head).
# Sized between body (19pt) and title (36pt) so the column header reads as a
# proper sub-title — "本文が2列なのは、それぞれにサブタイトルがあるから".
TWOCOL_HEAD_PT = 28


from lxml import etree


def _strip_bullet(paragraph):
    """段落のpPrからブレットを除去し、ハンギングインデントもクリアする。
    サブタイトル行（TwoColのhead）に使う。"""
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buFont", "a:buClr", "a:buSzPct", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    etree.SubElement(pPr, qn("a:buNone"))
    pPr.set("indent", "0")
    pPr.set("marL", "0")


class Deck:
    def __init__(self, template: str = None, include_intro: bool = True):
        """Michikusaテンプレをロードしてデッキを作る。

        include_intro=True (デフォルト): 冒頭にMichikusa会社紹介スライド5枚
        (Section/Bullets ×4) を自動で含む `michikusa_intro.pptx` を使う。
        生成するスライドはこれらの後ろに追加される。
        include_intro=False: 空テンプレの `michikusa.pptx` を使う。
        """
        if template:
            path = template
            preserve = 0
        elif include_intro and os.path.exists(_TEMPLATE_INTRO):
            path = _TEMPLATE_INTRO
            preserve = -1  # イントロpptx内のスライドは全て保持
        else:
            path = _TEMPLATE
            preserve = 0
        self.prs = Presentation(path)
        self._layouts = {l.name: l for l in self.prs.slide_master.slide_layouts}
        assert {"Cover", "Section", "Bullets", "TwoCol"}.issubset(self._layouts), \
            f"Template missing layouts: {self._layouts.keys()}"
        if preserve < 0:
            preserve = len(self.prs.slides)
        self._clear(preserve=preserve)

    def _clear(self, preserve: int = 0):
        """スライドを削除する。preserve 件は先頭から保持する。"""
        sldIdLst = self.prs.slides._sldIdLst
        slides = list(sldIdLst)
        for s in slides[preserve:]:
            rId = s.get(qn("r:id"))
            self.prs.part.drop_rel(rId)
            sldIdLst.remove(s)

    @staticmethod
    def _ph(slide, idx):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    @staticmethod
    def _set(ph, text):
        tf = ph.text_frame
        lines = text.split("\n")
        tf.text = lines[0]
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line

    def cover(self, title: str, subtitle: str = ""):
        s = self.prs.slides.add_slide(self._layouts["Cover"])
        self._set(self._ph(s, 0), jp_wrap(title, WRAP["cover_title"]))
        if subtitle:
            self._set(self._ph(s, 1), jp_wrap(subtitle, WRAP["cover_subtitle"]))
        return s

    def section(self, title: str):
        s = self.prs.slides.add_slide(self._layouts["Section"])
        self._set(self._ph(s, 0), jp_wrap(title, WRAP["section"]))
        return s

    def bullets(self, title: str, items: list):
        s = self.prs.slides.add_slide(self._layouts["Bullets"])
        self._set(self._ph(s, 0), jp_wrap(title, WRAP["bullets_title"]))
        body = "\n".join(jp_wrap(it, WRAP["bullets_body"]) for it in items)
        self._set(self._ph(s, 1), body)
        return s

    def twocol(self, title: str, left: list, right: list,
               left_head: str = None, right_head: str = None):
        """2列スライドを追加。left_head/right_head は専用サブタイトル
        プレースホルダー(idx=10/11) に入る。旧テンプレートでは本文の先頭段落に
        フォールバック。"""
        s = self.prs.slides.add_slide(self._layouts["TwoCol"])
        self._set(self._ph(s, 0), jp_wrap(title, WRAP["twocol_title"]))

        left_sub = self._ph(s, 10)
        right_sub = self._ph(s, 11)
        has_sub_ph = left_sub is not None and right_sub is not None

        if has_sub_ph:
            if left_head:
                self._set(left_sub, jp_wrap(left_head, WRAP["twocol_head"]))
            if right_head:
                self._set(right_sub, jp_wrap(right_head, WRAP["twocol_head"]))
            self._fill_col_body(self._ph(s, 1), left)
            self._fill_col_body(self._ph(s, 2), right)
        else:
            # 旧テンプレ: 本文プレースホルダー先頭段落にサブタイトルを載せる
            self._fill_col(self._ph(s, 1), left_head, left)
            self._fill_col(self._ph(s, 2), right_head, right)
        return s

    @staticmethod
    def _fill_col_body(ph, items):
        """TwoColの本文カラム(Left/Right Column)を箇条書きで埋める。"""
        tf = ph.text_frame
        body_lines = []
        for it in items:
            body_lines.extend(jp_wrap(it, WRAP["twocol_body"]).split("\n"))
        tf.text = body_lines[0] if body_lines else ""
        for line in body_lines[1:]:
            p = tf.add_paragraph()
            p.text = line

    @staticmethod
    def _fill_col(ph, head, items):
        """[旧テンプレ用フォールバック] TwoColの1カラムを埋める。
        先頭にheadがあればサブタイトルとして本文より大きく(TWOCOL_HEAD_PT)
        ブレットなしでレンダリングする。"""
        tf = ph.text_frame
        body_lines = []
        for it in items:
            body_lines.extend(jp_wrap(it, WRAP["twocol_body"]).split("\n"))

        if head:
            head_lines = jp_wrap(head, WRAP["twocol_head"]).split("\n")
            tf.text = head_lines[0]
            _strip_bullet(tf.paragraphs[0])
            for r in tf.paragraphs[0].runs:
                r.font.size = Pt(TWOCOL_HEAD_PT)
            for line in head_lines[1:]:
                p = tf.add_paragraph()
                p.text = line
                _strip_bullet(p)
                for r in p.runs:
                    r.font.size = Pt(TWOCOL_HEAD_PT)
            for line in body_lines:
                p = tf.add_paragraph()
                p.text = line
        else:
            tf.text = body_lines[0] if body_lines else ""
            for line in body_lines[1:]:
                p = tf.add_paragraph()
                p.text = line

    def save(self, path: str):
        self.prs.save(path)
        print(f"Saved: {path} ({len(self.prs.slides)} slides)")
        return path
